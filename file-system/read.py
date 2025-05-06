import sys
import os
import csv
import io
import json
import time
from typing import Dict, Any, Union, List, Optional

# For PDF processing
try:
    import PyPDF2
except ImportError:
    print("PyPDF2 not installed. To read PDF files: pip install PyPDF2")

# For enhanced PDF processing with better image detection
try:
    import fitz  # PyMuPDF
    has_pymupdf = True
except ImportError:
    has_pymupdf = False
    print("PyMuPDF not installed. For better PDF image detection: pip install pymupdf")

# For PDF table extraction
try:
    import tabula
    has_tabula = True
except ImportError:
    has_tabula = False
    print("Tabula-py not installed. For better PDF table extraction: pip install tabula-py")

# For PDF text extraction with layout preservation
try:
    import pdfplumber
    has_pdfplumber = True
except ImportError:
    has_pdfplumber = False
    print("pdfplumber not installed. For better PDF text extraction: pip install pdfplumber")

# For EPUB e-books
try:
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup
    has_epub_support = True
except ImportError:
    has_epub_support = False
    print("EPUB support not available. To read EPUB files: pip install ebooklib beautifulsoup4")

# For Microsoft Office documents
try:
    import docx  # for .docx files
    from docx.oxml.ns import qn  # for accessing embedded objects
    from docx.oxml import OxmlElement
except ImportError:
    print("python-docx not installed. To read Word files: pip install python-docx")

try:
    import openpyxl  # for .xlsx files
    from openpyxl.drawing.image import Image as XlsxImage
except ImportError:
    print("openpyxl not installed. To read Excel files: pip install openpyxl")

try:
    from pptx import Presentation  # for .pptx files
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("python-pptx not installed. To read PowerPoint files: pip install python-pptx")

# For image detection and properties
try:
    from PIL import Image
    has_pil = True
except ImportError:
    print("Pillow not installed. For better image analysis: pip install Pillow")
    has_pil = False

# For RTF documents
try:
    import striprtf.striprtf as striprtf
    has_rtf_support = True
except ImportError:
    has_rtf_support = False
    print("RTF support not available. To read RTF files: pip install striprtf")

def read_text_file(file_path: str) -> str:
    """Read content from text-based files with proper encoding detection."""
    encodings = ['utf-8', 'latin-1', 'windows-1252', 'ascii']
    
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            continue

    # If all encodings fail, try binary mode as last resort
    try:
        with open(file_path, 'rb') as f:
            return f.read().decode('utf-8', errors='replace')
    except Exception as e:
        return f"Error: Could not decode file with available encodings: {str(e)}"

def read_pdf_file(file_path: str) -> str:
    """Extract text and identify images from PDF files with layout preservation."""
    content = []
    
    # Try pdfplumber first if available (better text layout preservation)
    if has_pdfplumber:
        try:
            content.append("--- Document Metadata ---")
            pdf = pdfplumber.open(file_path)
            
            # Basic document info
            content.append(f"PDF Document: {os.path.basename(file_path)}")
            content.append(f"Number of pages: {len(pdf.pages)}")
            
            # Try to extract metadata
            if hasattr(pdf, 'metadata') and pdf.metadata:
                for key, value in pdf.metadata.items():
                    if value and str(value).strip():
                        # Clean up key name
                        clean_key = key[1:] if isinstance(key, str) and key.startswith('/') else key
                        content.append(f"{clean_key}: {value}")
            
            content.append("-" * 40)
            
            # Process each page
            for i, page in enumerate(pdf.pages):
                content.append(f"--- Page {i + 1} ---")
                
                # Extract page dimensions
                width, height = page.width, page.height
                content.append(f"Page dimensions: {width:.2f} x {height:.2f} points")
                
                # Extract text with layout preservation
                page_text = page.extract_text(x_tolerance=3, y_tolerance=3)
                
                # Try to detect tables
                tables = page.extract_tables()
                if tables:
                    content.append(f"[Contains {len(tables)} table{'s' if len(tables) > 1 else ''}]")
                    
                    # Format each table
                    for t_idx, table in enumerate(tables):
                        content.append(f"\nTable {t_idx + 1}:")
                        
                        # Format the table with proper alignment
                        for row in table:
                            # Clean row data and handle None values
                            cleaned_row = [str(cell).strip() if cell is not None else "" for cell in row]
                            row_text = " | ".join(cleaned_row)
                            content.append(row_text)
                        
                        content.append("")  # Add spacing after table
                
                # Extract images if available
                try:
                    images = page.images
                    if images:
                        content.append(f"[Contains {len(images)} image{'s' if len(images) > 1 else ''}]")
                        for img_idx, img in enumerate(images):
                            if 'width' in img and 'height' in img:
                                content.append(f"  Image {img_idx+1}: {img.get('width', '?')}x{img.get('height', '?')} pixels")
                            else:
                                content.append(f"  Image {img_idx+1}")
                except Exception as e:
                    pass  # Skip if image extraction fails
                
                # Add the page text with layout preservation
                if not page_text or page_text.isspace():
                    content.append("[This page appears to be empty or contains only non-text elements]")
                else:
                    content.append(page_text)
            
            pdf.close()
            return "\n\n".join(content)
        except Exception as e:
            print(f"pdfplumber processing failed: {str(e)}. Trying PyMuPDF.")
    
    # Try PyMuPDF if available
    if has_pymupdf:
        try:
            pdf_document = fitz.open(file_path)
            
            # Extract document metadata
            content.append("--- Document Metadata ---")
            metadata = pdf_document.metadata
            for key, value in metadata.items():
                if value:
                    content.append(f"{key}: {value}")
            content.append("-" * 40)
            
            # Document summary
            content.append(f"PDF Document: {os.path.basename(file_path)}")
            content.append(f"Number of pages: {len(pdf_document)}")
            content.append(f"PDF Version: {pdf_document.pdf_version}")
            if pdf_document.is_encrypted:
                content.append("Status: Encrypted")
            
            content.append("-" * 40)
            
            # Try using tabula for table extraction if available
            tables_by_page = {}
            if has_tabula:
                try:
                    # Extract tables from all pages
                    all_tables = tabula.read_pdf(file_path, pages='all', multiple_tables=True)
                    
                    # Organize tables by page
                    for idx, table in enumerate(all_tables):
                        # Tabula doesn't provide page numbers directly, so we estimate
                        # This is a limitation, might need adjustment based on PDF structure
                        page_num = idx % len(pdf_document) + 1  # Basic estimation
                        if page_num not in tables_by_page:
                            tables_by_page[page_num] = []
                        tables_by_page[page_num].append(table)
                except Exception as e:
                    print(f"Tabula table extraction failed: {str(e)}")
            
            # Process each page
            for page_num, page in enumerate(pdf_document):
                content.append(f"--- Page {page_num + 1} ---")
                
                # Get page size
                width, height = page.rect.width, page.rect.height
                content.append(f"Page dimensions: {width:.2f} x {height:.2f} points")
                
                # We'll extract content in blocks to preserve layout and sequence
                # Get the page elements in order
                page_elements = []
                
                # First, extract text blocks with position information
                # Sort by y position first (top to bottom), then x position (left to right)
                blocks = page.get_text("dict", sort=True)["blocks"]
                
                # Track text blocks with their position
                text_blocks = []
                for block in blocks:
                    if block["type"] == 0:  # Text block
                        text = ""
                        for line in block["lines"]:
                            for span in line["spans"]:
                                text += span["text"]
                            text += "\n"
                        
                        # Store the text block with its position
                        if text.strip():
                            y_pos = block["bbox"][1]  # Top y-coordinate
                            x_pos = block["bbox"][0]  # Left x-coordinate
                            text_blocks.append({
                                "type": "text",
                                "y_pos": y_pos,
                                "x_pos": x_pos,
                                "text": text.strip()
                            })
                
                # Get image information with position
                image_blocks = []
                for img_index, img in enumerate(page.get_images(full=True)):
                    xref = img[0]  # image reference number
                    
                    # Get image details
                    try:
                        base_image = pdf_document.extract_image(xref)
                        if base_image:
                            # Try to get image position on page
                            rect = None
                            for item in page.get_drawings():
                                if item["type"] == "image" and xref == item.get("xref"):
                                    rect = item["rect"]
                                    break
                            
                            if rect:
                                y_pos = rect[1]  # Top y-coordinate
                                x_pos = rect[0]  # Left x-coordinate
                                img_ext = base_image.get("ext", "")
                                img_width = base_image.get("width", "unknown")
                                img_height = base_image.get("height", "unknown")
                                
                                image_blocks.append({
                                    "type": "image",
                                    "y_pos": y_pos,
                                    "x_pos": x_pos,
                                    "desc": f"Image {img_index+1}: {img_width}x{img_height} {img_ext}"
                                })
                    except:
                        # If we can't get position, add it to the end
                        image_blocks.append({
                            "type": "image",
                            "y_pos": height,  # Put at bottom if position unknown
                            "x_pos": 0,
                            "desc": f"Image {img_index+1}"
                        })
                
                # Merge all blocks and sort by position
                page_elements = text_blocks + image_blocks
                page_elements.sort(key=lambda x: (x["y_pos"], x["x_pos"]))
                
                # Check for tables from tabula
                table_blocks = []
                if page_num + 1 in tables_by_page:
                    tables = tables_by_page[page_num + 1]
                    for t_idx, table in enumerate(tables):
                        # We don't have position info from tabula, so estimate
                        # based on where tables typically appear in the document
                        # This is approximate and may not perfectly match the layout
                        est_y_pos = height * 0.3 * (t_idx + 1)  # Estimate
                        
                        table_blocks.append({
                            "type": "table",
                            "y_pos": est_y_pos,
                            "x_pos": width / 4,  # Center-left of page
                            "table": table
                        })
                
                # Add tables to page elements if we have any
                if table_blocks:
                    page_elements.extend(table_blocks)
                    # Re-sort with tables included
                    page_elements.sort(key=lambda x: (x["y_pos"], x["x_pos"]))
                
                # If we have no elements, check if page has text using plain extraction
                if not page_elements:
                    page_text = page.get_text()
                    if page_text.strip():
                        content.append(page_text.strip())
                    else:
                        content.append("[This page appears to be empty or contains only non-text elements]")
                else:
                    # Output elements in order
                    for element in page_elements:
                        if element["type"] == "text":
                            content.append(element["text"])
                        elif element["type"] == "image":
                            content.append(f"[{element['desc']}]")
                        elif element["type"] == "table":
                            content.append("\n--- Table ---")
                            table_str = element["table"].to_string(index=False)
                            content.append(table_str)
                            content.append("--- End Table ---")
                
                # Check for links and annotations at the end
                links = page.get_links()
                if links:
                    content.append(f"[Contains {len(links)} link{'s' if len(links) > 1 else ''}]")
                
                annots = page.annots()
                if annots:
                    content.append(f"[Contains {len(annots)} annotation{'s' if len(annots) > 1 else ''}]")
            
            return "\n\n".join(content)
        except Exception as e:
            print(f"PyMuPDF processing failed: {str(e)}. Falling back to PyPDF2.")
    
    # Use PyPDF2 as final fallback
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            num_pages = len(pdf_reader.pages)
            
            # Try to get document info
            content.append("--- Document Metadata ---")
            info = pdf_reader.metadata
            if info:
                for key in info:
                    if info.get(key):
                        # Clean up key name by removing leading slash
                        clean_key = key[1:] if key.startswith('/') else key
                        content.append(f"{clean_key}: {info.get(key)}")
            
            # Add document summary
            content.append(f"PDF Document: {os.path.basename(file_path)}")
            content.append(f"Number of pages: {num_pages}")
            content.append("-" * 40)
            
            # Extract text from each page
            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                
                # Start page content
                content.append(f"--- Page {page_num + 1} ---")
                
                # Extract text with better layout handling
                try:
                    page_text = page.extract_text(extraction_mode="layout")
                except:
                    # Fallback to regular extraction if layout mode is not available
                    page_text = page.extract_text()
                
                # Try to detect images and objects
                try:
                    if '/Resources' in page:
                        resources = page['/Resources']
                        
                        # Check for XObject resources (images, forms, etc.)
                        if '/XObject' in resources:
                            xobjects = resources['/XObject']
                            if isinstance(xobjects, dict):
                                image_count = 0
                                form_count = 0
                                
                                for obj_name, obj_ref in xobjects.items():
                                    if '/Subtype' in obj_ref and obj_ref['/Subtype'] == '/Image':
                                        image_count += 1
                                    elif '/Subtype' in obj_ref and obj_ref['/Subtype'] == '/Form':
                                        form_count += 1
                                
                                if image_count > 0:
                                    content.append(f"[Contains {image_count} image{'s' if image_count > 1 else ''}]")
                                if form_count > 0:
                                    content.append(f"[Contains {form_count} form object{'s' if form_count > 1 else ''}]")
                except:
                    # If resource detection fails, just continue
                    pass
                
                # Handle empty pages
                if not page_text or page_text.isspace():
                    content.append("[This page appears to be empty or contains only non-text elements]")
                else:
                    # Format text with proper paragraph separation
                    paragraphs = page_text.split('\n\n')
                    formatted_text = []
                    for para in paragraphs:
                        # Ensure there are no artificial line breaks within paragraphs
                        # but preserve real paragraph breaks
                        lines = para.split('\n')
                        formatted_para = ' '.join(line.strip() for line in lines if line.strip())
                        if formatted_para:
                            formatted_text.append(formatted_para)
                    
                    content.append('\n\n'.join(formatted_text))
        
        return "\n\n".join(content)
    except Exception as e:
        return f"Error reading PDF file: {str(e)}"

def read_docx_file(file_path: str) -> str:
    """Extract text and identify objects from Microsoft Word (.docx) files in sequential order."""
    try:
        doc = docx.Document(file_path)
        full_text = []
        
        # Document properties first
        try:
            full_text.append("--- Document Properties ---")
            core_props = doc.core_properties
            if hasattr(core_props, 'title') and core_props.title:
                full_text.append(f"Title: {core_props.title}")
            if hasattr(core_props, 'author') and core_props.author:
                full_text.append(f"Author: {core_props.author}")
            if hasattr(core_props, 'created') and core_props.created:
                full_text.append(f"Created: {core_props.created}")
            if hasattr(core_props, 'modified') and core_props.modified:
                full_text.append(f"Modified: {core_props.modified}")
            if hasattr(core_props, 'comments') and core_props.comments:
                full_text.append(f"Comments: {core_props.comments}")
            if hasattr(core_props, 'category') and core_props.category:
                full_text.append(f"Category: {core_props.category}")
            if hasattr(core_props, 'subject') and core_props.subject:
                full_text.append(f"Subject: {core_props.subject}")
            if hasattr(core_props, 'keywords') and core_props.keywords:
                full_text.append(f"Keywords: {core_props.keywords}")
            full_text.append("-" * 40)
        except:
            pass  # Ignore if properties can't be accessed
        
        # Document statistics
        full_text.append("--- Document Statistics ---")
        full_text.append(f"Paragraphs: {len(doc.paragraphs)}")
        full_text.append(f"Sections: {len(doc.sections)}")
        full_text.append(f"Tables: {len(doc.tables)}")
        
        # Detect images and other embedded objects
        try:
            # Scan for embedded objects by looking at the underlying XML
            image_count = 0
            chart_count = 0
            shape_count = 0
            drawing_count = 0
            
            for rel in doc.part.rels.values():
                if rel.reltype == 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image':
                    image_count += 1
                elif rel.reltype == 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart':
                    chart_count += 1
                elif rel.reltype == 'http://schemas.microsoft.com/office/2007/relationships/shape':
                    shape_count += 1
                elif 'drawing' in rel.reltype:
                    drawing_count += 1
            
            if image_count > 0:
                full_text.append(f"Images: {image_count}")
                
                # Get more detailed image information if possible
                try:
                    image_details = []
                    for rel_id, rel in doc.part.rels.items():
                        if rel.reltype == 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image':
                            try:
                                # Try to get image dimensions
                                image_part = rel.target_part
                                if has_pil and image_part and hasattr(image_part, 'blob'):
                                    from io import BytesIO
                                    img = Image.open(BytesIO(image_part.blob))
                                    image_details.append(f"  Image {rel_id}: {img.format} {img.width}x{img.height}")
                                else:
                                    image_details.append(f"  Image {rel_id}")
                            except:
                                image_details.append(f"  Image {rel_id}")
                    
                    if image_details:
                        full_text.append("Image details:")
                        full_text.extend(image_details)
                except:
                    pass  # Skip if detailed image info can't be obtained
            
            if chart_count > 0:
                full_text.append(f"Charts: {chart_count}")
            if shape_count > 0:
                full_text.append(f"Shapes: {shape_count}")
            if drawing_count > 0:
                full_text.append(f"Drawings: {drawing_count}")
            
            full_text.append("-" * 40)
        except:
            full_text.append("-" * 40)  # Add separator even if object detection fails
        
        # Main content extraction - preserving sequence
        full_text.append("--- Document Content (in sequential order) ---")
        
        # We need to track all content elements (paragraphs, tables, images) and their positions
        content_elements = []
        
        # Track paragraphs and their sequence
        for i, para in enumerate(doc.paragraphs):
            # Skip empty paragraphs
            if not para.text.strip():
                continue
                
            # Check for headings
            heading_level = 0
            if para.style and para.style.name.startswith('Heading'):
                try:
                    heading_level = int(para.style.name.replace('Heading', ''))
                except ValueError:
                    heading_level = 0
            
            # Check for images or other objects within this paragraph
            has_objects = False
            has_drawings = False
            drawing_texts = []
            
            for run in para.runs:
                # Check for embedded images
                if run.element.findall('.//'+qn('w:drawing')) or run.element.findall('.//'+qn('w:pict')):
                    has_objects = True
                    
                    # Try to extract text from drawing elements (text boxes, etc.)
                    for drawing in run.element.findall('.//'+qn('w:drawing')):
                        has_drawings = True
                        # Try to extract any text from the drawing
                        try:
                            # Look for text elements in the drawing
                            texts = drawing.findall('.//' + qn('w:t'))
                            if texts:
                                drawing_text = ' '.join([t.text for t in texts if t.text])
                                if drawing_text.strip():
                                    drawing_texts.append(drawing_text.strip())
                        except:
                            pass
            
            # Create the paragraph text with appropriate heading level
            if heading_level > 0:
                para_text = f"{'#' * heading_level} {para.text}"
            else:
                para_text = para.text
            
            # Add information about embedded objects/drawings
            if has_objects:
                if drawing_texts:
                    para_text += f" [Contains drawing/image with text: {'; '.join(drawing_texts)}]"
                else:
                    para_text += " [Contains embedded object(s)]"
            
            # Store the paragraph with its document position
            content_elements.append({
                'type': 'paragraph',
                'position': i,
                'content': para_text,
                'has_objects': has_objects,
                'has_drawings': has_drawings,
                'drawing_texts': drawing_texts
            })
        
        # Track tables and their sequence
        for i, table in enumerate(doc.tables):
            # Find the paragraph immediately before the table to determine position
            table_position = -1
            
            # This is an approximation as the exact position might be harder to determine
            # We'll use XML to try to find where this table appears
            try:
                # Get all paragraphs and tables in order from the body
                body = doc._body._body
                elements = body.findall('.//')
                
                # Walk through elements to find this table's position
                current_position = 0
                for elem in elements:
                    if elem.tag.endswith('tbl'):  # Table element
                        if current_position == i:  # If this is our target table
                            table_position = current_position
                            break
                        current_position += 1
                    elif elem.tag.endswith('p'):  # Paragraph element
                        current_position += 1
            except:
                # If we can't determine position, just use the table index
                table_position = i * 1000  # Arbitrary large value to put tables after paragraphs
            
            # Format table data
            table_data = []
            col_widths = []
            
            # Process the table - similar to previous implementation
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    # Combine all text in the cell, handling paragraphs
                    cell_text = '\n'.join([p.text for p in cell.paragraphs if p.text.strip()])
                    
                    # Check for embedded objects in the cell
                    has_objects = False
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            if run.element.findall('.//'+qn('w:drawing')) or run.element.findall('.//'+qn('w:pict')):
                                has_objects = True
                                break
                    
                    if has_objects:
                        cell_text += " [Contains embedded object(s)]"
                        
                    # Clean up the cell text
                    cell_text = cell_text.strip()
                    row_data.append(cell_text)
                
                # Update column widths
                for i, cell_text in enumerate(row_data):
                    # For multiline content, get the max line width
                    lines = cell_text.split('\n')
                    max_line_width = max([len(line) for line in lines]) if lines else 0
                    
                    while len(col_widths) <= i:
                        col_widths.append(0)
                    col_widths[i] = max(col_widths[i], max_line_width)
                
                table_data.append(row_data)
            
            # Store the formatted table data
            content_elements.append({
                'type': 'table',
                'position': table_position,
                'content': table_data,
                'col_widths': col_widths
            })
        
        # Sort content elements by their position to maintain document flow
        content_elements.sort(key=lambda x: x['position'])
        
        # Now process each element in order
        for element in content_elements:
            if element['type'] == 'paragraph':
                full_text.append(element['content'])
            elif element['type'] == 'table':
                table_data = element['content']
                col_widths = element['col_widths']
                
                full_text.append("\n--- Table ---")
                
                # Format and output the table
                if table_data:
                    # Header row
                    header_row = table_data[0]
                    header_line = []
                    for i, cell in enumerate(header_row):
                        width = min(col_widths[i], 30)  # Limit width to reasonable size
                        header_line.append(cell.ljust(width))
                    full_text.append("| " + " | ".join(header_line) + " |")
                    
                    # Separator row
                    separator = []
                    for i in range(len(header_row)):
                        width = min(col_widths[i], 30)
                        separator.append("-" * width)
                    full_text.append("| " + " | ".join(separator) + " |")
                    
                    # Data rows
                    for row_idx, row_data in enumerate(table_data):
                        if row_idx == 0:  # Skip header row as we already displayed it
                            continue
                        
                        # Handle multiline cells
                        row_lines = []
                        max_lines = 1
                        
                        for i, cell in enumerate(row_data):
                            if not cell:
                                row_lines.append([""])
                                continue
                                
                            lines = cell.split('\n')
                            row_lines.append(lines)
                            max_lines = max(max_lines, len(lines))
                        
                        # Now create and output each line of the row
                        for line_idx in range(max_lines):
                            line_parts = []
                            for col_idx, cell_lines in enumerate(row_lines):
                                if line_idx < len(cell_lines):
                                    content = cell_lines[line_idx]
                                else:
                                    content = ""
                                    
                                width = min(col_widths[col_idx], 30)
                                line_parts.append(content.ljust(width))
                            
                            full_text.append("| " + " | ".join(line_parts) + " |")
                
                full_text.append("--- End Table ---\n")
        
        # Check for image-only or drawing-only objects not attached to paragraphs
        # This requires special handling at the XML level
        try:
            body = doc._body._body
            
            # Function to extract text from a drawing object
            def extract_drawing_text(drawing_elem):
                drawing_texts = []
                try:
                    texts = drawing_elem.findall('.//' + qn('w:t'))
                    if texts:
                        drawing_text = ' '.join([t.text for t in texts if t.text])
                        if drawing_text.strip():
                            drawing_texts.append(drawing_text.strip())
                except:
                    pass
                return drawing_texts
            
            # Find standalone drawing objects
            standalone_drawings = []
            for elem in body.findall('.//'+qn('w:drawing')):
                # Check if this drawing is directly in the body, not inside a paragraph we already processed
                parent = elem.getparent()
                if parent is not None and parent.tag.endswith('r'):  # Inside a run
                    run_parent = parent.getparent()
                    if run_parent is not None and run_parent.tag.endswith('p'):  # Inside a paragraph
                        # This might be a standalone drawing, check if we've already captured it
                        # This logic is approximate - ideally we'd match exact elements
                        drawing_texts = extract_drawing_text(elem)
                        if drawing_texts:
                            # Check if this text is in any of our existing paragraphs with drawings
                            found = False
                            for e in content_elements:
                                if e['type'] == 'paragraph' and e['has_drawings']:
                                    for dt in e['drawing_texts']:
                                        if dt in drawing_texts:
                                            found = True
                                            break
                            
                            if not found:
                                standalone_drawings.append({
                                    'type': 'drawing',
                                    'texts': drawing_texts
                                })
            
            # Add any standalone drawings we found
            if standalone_drawings:
                full_text.append("\n--- Standalone Drawings/Text Boxes ---")
                for i, drawing in enumerate(standalone_drawings):
                    full_text.append(f"Drawing/Text Box {i+1}:")
                    for text in drawing['texts']:
                        full_text.append(f"  {text}")
                full_text.append("--- End Drawings ---\n")
                
        except:
            pass  # Skip if we can't process standalone drawings
        
        return "\n".join(full_text)
    except Exception as e:
        return f"Error reading DOCX file: {str(e)}"

def read_xlsx_file(file_path: str) -> str:
    """Extract data and identify objects from Microsoft Excel (.xlsx) files."""
    try:
        workbook = openpyxl.load_workbook(file_path, data_only=True)  # data_only=True to get values instead of formulas
        result = []
        
        # Workbook properties
        result.append("--- Workbook Properties ---")
        result.append(f"Filename: {os.path.basename(file_path)}")
        result.append(f"Number of sheets: {len(workbook.sheetnames)}")
        
        # Try to get document properties
        try:
            props = workbook.properties
            if props.title:
                result.append(f"Title: {props.title}")
            if props.subject:
                result.append(f"Subject: {props.subject}")
            if props.creator:
                result.append(f"Creator: {props.creator}")
            if props.created:
                result.append(f"Created: {props.created}")
        except:
            pass  # Skip if properties aren't accessible
            
        result.append("-" * 40)
        
        # Process each worksheet
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            result.append(f"--- Sheet: {sheet_name} ---")
            
            # Get sheet dimensions and properties
            min_row, min_col, max_row, max_col = 1, 1, sheet.max_row, sheet.max_column
            result.append(f"Dimensions: {min_col}:{min_row} to {max_col}:{max_row}")
            
            # Find charts
            chart_count = 0
            try:
                for chart_ref in sheet._charts:
                    chart_count += 1
                
                if chart_count > 0:
                    result.append(f"Charts: {chart_count}")
            except:
                pass
            
            # Find images with more details
            image_count = 0
            try:
                # Load workbook again without data_only to access images
                image_workbook = openpyxl.load_workbook(file_path, data_only=False)
                image_sheet = image_workbook[sheet_name]
                
                for image in image_sheet._images:
                    image_count += 1
                    
                if image_count > 0:
                    result.append(f"Images: {image_count}")
                    
                    # Try to get more details about images
                    image_details = []
                    for i, img in enumerate(image_sheet._images):
                        try:
                            if hasattr(img, 'width') and hasattr(img, 'height'):
                                image_details.append(f"  Image {i+1}: {img.width}x{img.height}")
                            else:
                                image_details.append(f"  Image {i+1}")
                        except:
                            image_details.append(f"  Image {i+1}")
                    
                    if image_details:
                        result.extend(image_details)
            except:
                pass  # Skip if images can't be accessed
                
            # Find merged cells
            if sheet.merged_cells:
                result.append(f"Merged cell ranges: {len(sheet.merged_cells.ranges)}")
            
            # Find conditional formatting
            if hasattr(sheet, 'conditional_formatting') and sheet.conditional_formatting:
                result.append(f"Conditional formatting rules: {len(sheet.conditional_formatting._cf_rules)}")
            
            result.append("-" * 40)
            
            # Skip completely empty sheets
            if max_row < 2 and max_col < 2:
                result.append("[This sheet appears to be empty]")
                continue
                
            # Process data as a proper table
            # First, determine column widths
            col_widths = {}
            data_rows = []
            
            # Reasonable limits for large sheets
            max_rows_to_process = min(1000, max_row)
            max_cols_to_process = min(50, max_col)
            
            # First pass: collect data and determine column widths
            for row_idx in range(min_row, max_rows_to_process + 1):
                row_data = []
                
                for col_idx in range(min_col, max_cols_to_process + 1):
                    cell = sheet.cell(row=row_idx, column=col_idx)
                    value = cell.value
                    
                    # Format the cell value
                    if value is None:
                        formatted_value = ""
                    elif isinstance(value, (int, float)):
                        # Handle numeric formatting
                        if cell.number_format and "0.00" in cell.number_format:
                            formatted_value = f"{value:.2f}"
                        elif cell.number_format and "0.0" in cell.number_format:
                            formatted_value = f"{value:.1f}"
                        elif cell.number_format and "#,##0" in cell.number_format:
                            formatted_value = f"{value:,}"
                        elif cell.number_format and "0%" in cell.number_format:
                            formatted_value = f"{value:.0%}"
                        elif cell.number_format and "0.00%" in cell.number_format:
                            formatted_value = f"{value:.2%}"
                        else:
                            formatted_value = str(value)
                    else:
                        formatted_value = str(value)
                    
                    row_data.append(formatted_value)
                    
                    # Update column width
                    width = len(formatted_value)
                    if col_idx not in col_widths or width > col_widths[col_idx]:
                        col_widths[col_idx] = min(width, 30)  # Cap width at 30 to avoid extremely wide columns
                
                data_rows.append(row_data)
            
            # Now output formatted table
            # Check if we have any data
            if not data_rows:
                result.append("[This sheet contains no data]")
                continue
                
            # Determine if first row is a header (simple heuristic)
            has_header = False
            if len(data_rows) > 1:
                first_row = data_rows[0]
                # Check if first row has more text cells while second row has more numeric cells
                text_cells_first_row = sum(1 for cell in first_row if cell and not cell.replace(".", "").isdigit())
                if text_cells_first_row > len(first_row) // 3:  # If at least 1/3 of cells are text
                    has_header = True
            
            # Now format the table with proper spacing
            if has_header and data_rows:
                # Header row
                header_parts = []
                for col_idx, value in enumerate(data_rows[0], start=min_col):
                    width = col_widths.get(col_idx, 10)
                    header_parts.append(value.ljust(width))
                result.append("| " + " | ".join(header_parts) + " |")
                
                # Separator line
                separator_parts = []
                for col_idx in range(min_col, min_col + len(data_rows[0])):
                    width = col_widths.get(col_idx, 10)
                    separator_parts.append("-" * width)
                result.append("| " + " | ".join(separator_parts) + " |")
                
                # Data rows
                for row_idx, row_data in enumerate(data_rows[1:], start=1):
                    row_parts = []
                    for col_idx, value in enumerate(row_data, start=min_col):
                        width = col_widths.get(col_idx, 10)
                        row_parts.append(value.ljust(width))
                    result.append("| " + " | ".join(row_parts) + " |")
            else:
                # No header, just output all rows uniformly
                for row_data in data_rows:
                    row_parts = []
                    for col_idx, value in enumerate(row_data, start=min_col):
                        width = col_widths.get(col_idx, 10)
                        row_parts.append(value.ljust(width))
                    result.append("| " + " | ".join(row_parts) + " |")
            
            # If we hit the row limit, indicate there's more data
            if max_row > max_rows_to_process:
                result.append(f"[... {max_row - max_rows_to_process} more rows not shown ...]")
            
            # Add space between sheets
            result.append("\n")
        
        return "\n".join(result)
    except Exception as e:
        return f"Error reading XLSX file: {str(e)}"

def read_pptx_file(file_path: str) -> str:
    """Extract text and identify objects from Microsoft PowerPoint (.pptx) files in sequential order."""
    try:
        prs = Presentation(file_path)
        full_text = []
        
        # Document properties
        full_text.append("--- Presentation Properties ---")
        full_text.append(f"Filename: {os.path.basename(file_path)}")
        full_text.append(f"Number of slides: {len(prs.slides)}")
        
        # Try to get core properties
        try:
            core_props = prs.core_properties
            if hasattr(core_props, 'title') and core_props.title:
                full_text.append(f"Title: {core_props.title}")
            if hasattr(core_props, 'author') and core_props.author:
                full_text.append(f"Author: {core_props.author}")
            if hasattr(core_props, 'subject') and core_props.subject:
                full_text.append(f"Subject: {core_props.subject}")
            if hasattr(core_props, 'keywords') and core_props.keywords:
                full_text.append(f"Keywords: {core_props.keywords}")
            if hasattr(core_props, 'created') and core_props.created:
                full_text.append(f"Created: {core_props.created}")
            if hasattr(core_props, 'modified') and core_props.modified:
                full_text.append(f"Modified: {core_props.modified}")
        except:
            pass  # Skip if properties can't be accessed
        
        full_text.append("-" * 40)
        
        # Slide master and layouts info
        try:
            slide_masters = prs.slide_masters
            full_text.append(f"Slide masters: {len(slide_masters)}")
            
            total_layouts = 0
            for master in slide_masters:
                total_layouts += len(master.slide_layouts)
            
            full_text.append(f"Slide layouts: {total_layouts}")
            full_text.append("-" * 40)
        except:
            full_text.append("-" * 40)
        
        # Process each slide
        for i, slide in enumerate(prs.slides):
            slide_content = [f"--- Slide {i+1} ---"]
            
            # Extract slide layout type if available
            try:
                layout_name = slide.slide_layout.name
                slide_content.append(f"Layout: {layout_name}")
            except:
                pass
            
            # Get slide dimensions
            try:
                if hasattr(prs, 'slide_width') and hasattr(prs, 'slide_height'):
                    slide_width = prs.slide_width / 914400  # Convert EMU to inches
                    slide_height = prs.slide_height / 914400
                    slide_content.append(f"Dimensions: {slide_width:.2f}\" × {slide_height:.2f}\"")
            except:
                pass
                
            # Extract text from slide title separately
            title_text = ""
            if slide.shapes.title:
                title_text = slide.shapes.title.text
                slide_content.append(f"Title: {title_text}")
            
            # Count shape types for summary
            shape_counts = {
                'text_boxes': 0,
                'pictures': 0,
                'charts': 0,
                'tables': 0,
                'diagrams': 0,
                'videos': 0,
                'other_shapes': 0
            }
            
            # Collect all shapes with their position information for sequential ordering
            slide_elements = []
            
            # Process all shapes to extract content and position
            for shape_idx, shape in enumerate(slide.shapes):
                # Skip the title as we've already handled it
                if shape == slide.shapes.title:
                    continue
                
                # Get shape position
                try:
                    top = shape.top / 914400  # EMU to inches
                    left = shape.left / 914400
                except:
                    # If we can't get position, use index as proxy for position
                    top = shape_idx * 10
                    left = 0
                
                # Check type and extract content
                if shape.has_text_frame:
                    if shape.text.strip():
                        shape_counts['text_boxes'] += 1
                        
                        # Store text with paragraph formatting preserved
                        shape_text = []
                        
                        # Preserve text formatting - check paragraphs and runs
                        for para in shape.text_frame.paragraphs:
                            para_text = []
                            for run in para.runs:
                                text = run.text.strip()
                                if text:
                                    # Check for formatting
                                    if run.font.bold:
                                        text = f"**{text}**"
                                    if run.font.italic:
                                        text = f"*{text}*"
                                    if run.font.underline:
                                        text = f"_{text}_"
                                    
                                    para_text.append(text)
                            
                            if para_text:
                                shape_text.append(" ".join(para_text))
                        
                        # Add the text content with formatting preserved
                        text_content = "\n".join(shape_text)
                        
                        slide_elements.append({
                            'type': 'text',
                            'top': top,
                            'left': left,
                            'content': text_content
                        })
                
                elif shape.has_table:
                    shape_counts['tables'] += 1
                    
                    # Process table data with column widths
                    table_data = []
                    col_widths = []
                    
                    # Extract table content
                    for row in shape.table.rows:
                        row_data = []
                        for cell in row.cells:
                            cell_text = cell.text.strip().replace('\n', ' ')
                            row_data.append(cell_text)
                        
                        # Update column widths
                        for i, text in enumerate(row_data):
                            while len(col_widths) <= i:
                                col_widths.append(0)
                            col_widths[i] = max(col_widths[i], len(text))
                        
                        table_data.append(row_data)
                    
                    slide_elements.append({
                        'type': 'table',
                        'top': top,
                        'left': left,
                        'content': table_data,
                        'col_widths': col_widths
                    })
                
                elif shape.has_chart:
                    shape_counts['charts'] += 1
                    
                    # Get chart type and data if possible
                    try:
                        chart_type = str(shape.chart.chart_type).split('.')[-1]
                        chart_desc = f"Chart: {chart_type}"
                    except:
                        chart_desc = "Chart"
                    
                    slide_elements.append({
                        'type': 'chart',
                        'top': top,
                        'left': left,
                        'content': chart_desc
                    })
                
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    shape_counts['pictures'] += 1
                    
                    # Get picture details if possible
                    try:
                        if hasattr(shape, 'width') and hasattr(shape, 'height'):
                            pic_width = shape.width / 914400  # EMU to inches
                            pic_height = shape.height / 914400
                            pic_desc = f"Picture: {pic_width:.1f}\" × {pic_height:.1f}\""
                        else:
                            pic_desc = "Picture"
                    except:
                        pic_desc = "Picture"
                    
                    slide_elements.append({
                        'type': 'picture',
                        'top': top,
                        'left': left,
                        'content': pic_desc
                    })
                
                elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                    shape_counts['videos'] += 1
                    
                    slide_elements.append({
                        'type': 'media',
                        'top': top,
                        'left': left,
                        'content': "Video/Media"
                    })
                
                elif shape.shape_type == MSO_SHAPE_TYPE.DIAGRAM:
                    shape_counts['diagrams'] += 1
                    
                    slide_elements.append({
                        'type': 'diagram',
                        'top': top,
                        'left': left,
                        'content': "Diagram/SmartArt"
                    })
                
                else:
                    shape_counts['other_shapes'] += 1
                    
                    # Try to get text from drawing objects or other shapes
                    shape_text = ""
                    try:
                        if hasattr(shape, 'text'):
                            shape_text = shape.text
                    except:
                        pass
                    
                    slide_elements.append({
                        'type': 'shape',
                        'top': top,
                        'left': left,
                        'content': f"Shape{': ' + shape_text if shape_text else ''}"
                    })
            
            # Add shape summary
            shape_summary = []
            for shape_type, count in shape_counts.items():
                if count > 0:
                    shape_summary.append(f"{count} {shape_type.replace('_', ' ')}")
            
            if shape_summary:
                slide_content.append("Objects: " + ", ".join(shape_summary))
            
            # Sort elements by position (top to bottom, then left to right)
            slide_elements.sort(key=lambda x: (x['top'], x['left']))
            
            # Now process all elements in sequential order
            slide_content.append("\n--- Slide Content (in sequential order) ---")
            
            for element in slide_elements:
                element_type = element['type']
                
                if element_type == 'text':
                    slide_content.append(f"Text Box:")
                    # Add indentation to preserve the text box structure
                    indented_text = "\n".join(f"  {line}" for line in element['content'].split('\n'))
                    slide_content.append(indented_text)
                
                elif element_type == 'picture':
                    slide_content.append(f"[{element['content']}]")
                
                elif element_type == 'chart':
                    slide_content.append(f"[{element['content']}]")
                
                elif element_type == 'diagram':
                    slide_content.append(f"[{element['content']}]")
                
                elif element_type == 'media':
                    slide_content.append(f"[{element['content']}]")
                
                elif element_type == 'shape':
                    slide_content.append(f"[{element['content']}]")
                
                elif element_type == 'table':
                    table_data = element['content']
                    col_widths = element['col_widths']
                    
                    slide_content.append("--- Table ---")
                    
                    # Format and output the table
                    if table_data:
                        # Header row
                        header_row = table_data[0]
                        header_parts = []
                        for i, text in enumerate(header_row):
                            width = min(col_widths[i], 30)  # Cap width at 30 chars
                            header_parts.append(text.ljust(width))
                        slide_content.append("| " + " | ".join(header_parts) + " |")
                        
                        # Separator
                        separator_parts = []
                        for i in range(len(header_row)):
                            width = min(col_widths[i], 30)
                            separator_parts.append("-" * width)
                        slide_content.append("| " + " | ".join(separator_parts) + " |")
                        
                        # Data rows
                        for row_idx, row_data in enumerate(table_data):
                            if row_idx == 0:  # Skip header as we already displayed it
                                continue
                                
                            row_parts = []
                            for i, text in enumerate(row_data):
                                width = min(col_widths[i], 30)
                                row_parts.append(text.ljust(width))
                            slide_content.append("| " + " | ".join(row_parts) + " |")
                    
                    slide_content.append("--- End Table ---")
            
            # Check if slide has minimal content
            if len(slide_content) <= 3:  # Only slide number, maybe layout and title
                slide_content.append("[This slide appears to have no text content or contains only non-text elements]")
            
            full_text.append("\n".join(slide_content))
        
        return "\n\n".join(full_text)
    except Exception as e:
        return f"Error reading PPTX file: {str(e)}"

def read_csv_file(file_path: str) -> str:
    """Read CSV files with enhanced delimiter detection and rich formatting."""
    try:
        # Try multiple encodings
        encodings = ['utf-8-sig', 'utf-8', 'latin-1', 'windows-1252']
        
        for encoding in encodings:
            try:
                # Try to detect delimiter and file structure
                with open(file_path, 'r', newline='', encoding=encoding) as csvfile:
                    # Read a sample to analyze
                    sample = csvfile.read(8192)  # Read larger sample for better detection
                    csvfile.seek(0)  # Reset to beginning of file
                    
                    # Analyze file structure
                    result = []
                    result.append(f"--- CSV File Analysis: {os.path.basename(file_path)} ---")
                    
                    # Count lines in the file
                    line_count = sample.count('\n')
                    if csvfile.read() != '':  # If there's more content after the sample
                        line_count = line_count + 1 + '...'
                    csvfile.seek(0)  # Reset again
                    
                    result.append(f"Estimated number of rows: {line_count if isinstance(line_count, int) else '1000+'}")
                    
                    # Try to detect the dialect
                    try:
                        dialect = csv.Sniffer().sniff(sample)
                        delimiter = dialect.delimiter
                        result.append(f"Detected delimiter: '{delimiter}'")
                        result.append(f"Quote character: '{dialect.quotechar}'")
                        
                        # Check if file has a header
                        has_header = csv.Sniffer().has_header(sample)
                        result.append(f"Has header row: {has_header}")
                    except:
                        # If sniffing fails, try common delimiters
                        for delimiter in [',', ';', '\t', '|']:
                            if delimiter in sample:
                                result.append(f"Using delimiter: '{delimiter}'")
                                break
                        has_header = True  # Assume header by default
                    
                    # Add separator
                    result.append("-" * 40)
                    
                    # Read and format CSV data
                    reader = csv.reader(csvfile, delimiter=delimiter)
                    rows = []
                    
                    # Read header if present
                    header = next(reader) if has_header else None
                    if header:
                        result.append(f"Column count: {len(header)}")
                        result.append(f"Headers: {', '.join(header)}")
                        result.append("-" * 40)
                        rows.append("\t".join(header))
                    
                    # Process data rows
                    row_count = 0
                    max_rows = 2000  # Reasonable limit for large files
                    column_stats = {}  # Track stats about each column
                    
                    for row in reader:
                        rows.append("\t".join(row))
                        
                        # Gather column statistics
                        for i, value in enumerate(row):
                            if i not in column_stats:
                                column_stats[i] = {'numeric': 0, 'empty': 0, 'total': 0}
                            
                            column_stats[i]['total'] += 1
                            if not value.strip():
                                column_stats[i]['empty'] += 1
                            elif value.strip().replace('.', '', 1).replace('-', '', 1).isdigit():
                                column_stats[i]['numeric'] += 1
                        
                        row_count += 1
                        if row_count >= max_rows:
                            rows.append("... (truncated due to large size)")
                            break
                    
                    # Add column analysis if we have enough data
                    if row_count > 10 and header:
                        result.append("\nColumn Analysis:")
                        for i, col_name in enumerate(header):
                            if i in column_stats and column_stats[i]['total'] > 0:
                                stats = column_stats[i]
                                pct_numeric = (stats['numeric'] / stats['total']) * 100
                                pct_empty = (stats['empty'] / stats['total']) * 100
                                
                                col_type = 'Numeric' if pct_numeric > 90 else 'Text'
                                if pct_empty > 50:
                                    col_type += " (Sparse)"
                                
                                result.append(f"  {col_name}: {col_type}")
                        
                        result.append("-" * 40)
                    
                    # Add the actual data
                    result.extend(rows)
                    
                    return "\n".join(result)
            except UnicodeDecodeError:
                continue  # Try next encoding
        
        # If all encodings fail, try as plain text
        return read_text_file(file_path)
    except Exception as e:
        # Fall back to simple text reading if CSV parsing fails
        try:
            return read_text_file(file_path)
        except Exception as nested_e:
            return f"Error reading CSV file: {str(e)}, Nested error: {str(nested_e)}"

def read_epub_file(file_path: str) -> str:
    """Extract content from EPUB e-books."""
    if not has_epub_support:
        return "EPUB support not available. Install required packages with: pip install ebooklib beautifulsoup4"
    
    try:
        content = []
        book = epub.read_epub(file_path)
        
        # Extract metadata
        content.append("--- EPUB Metadata ---")
        content.append(f"Title: {book.get_metadata('DC', 'title')[0][0] if book.get_metadata('DC', 'title') else 'Unknown'}")
        
        # Get author(s)
        authors = book.get_metadata('DC', 'creator')
        if authors:
            author_list = [author[0] for author in authors]
            content.append(f"Author(s): {', '.join(author_list)}")
        
        # Get language
        languages = book.get_metadata('DC', 'language')
        if languages:
            content.append(f"Language: {languages[0][0]}")
            
        # Get other metadata
        identifiers = book.get_metadata('DC', 'identifier')
        if identifiers:
            for identifier in identifiers:
                if identifier[1].get('id') == 'ISBN':
                    content.append(f"ISBN: {identifier[0]}")
        
        publishers = book.get_metadata('DC', 'publisher')
        if publishers:
            content.append(f"Publisher: {publishers[0][0]}")
            
        dates = book.get_metadata('DC', 'date')
        if dates:
            content.append(f"Date: {dates[0][0]}")
            
        content.append("-" * 40)
        
        # Get table of contents
        toc = book.toc
        if toc:
            content.append("--- Table of Contents ---")
            
            def process_toc_entries(entries, level=0):
                toc_content = []
                for entry in entries:
                    if isinstance(entry, tuple) and len(entry) >= 2:
                        title, href = entry[0], entry[1]
                        toc_content.append(f"{'  ' * level}• {title}")
                    elif isinstance(entry, list):
                        toc_content.extend(process_toc_entries(entry, level + 1))
                return toc_content
            
            content.extend(process_toc_entries(toc))
            content.append("-" * 40)
        
        # Count items and get document statistics
        content.append("--- Document Statistics ---")
        content.append(f"Spine items: {len(book.spine)}")
        content.append(f"Total items: {len(book.items)}")
        
        # Count images
        image_count = 0
        css_count = 0
        html_count = 0
        
        for item in book.items:
            if item.media_type and item.media_type.startswith('image/'):
                image_count += 1
            elif item.media_type == 'text/css':
                css_count += 1
            elif item.media_type == 'application/xhtml+xml':
                html_count += 1
        
        content.append(f"HTML documents: {html_count}")
        content.append(f"CSS stylesheets: {css_count}")
        content.append(f"Images: {image_count}")
        content.append("-" * 40)
        
        # Extract text content from HTML documents
        content.append("--- Content ---")
        
        # Helper function to extract text from HTML
        def chapter_to_text(html_content):
            try:
                soup = BeautifulSoup(html_content, 'html.parser')
                
                # Extract title if available
                title = soup.find('title')
                title_text = f"Chapter: {title.text}\n" if title else ""
                
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.extract()
                
                # Get text
                text = soup.get_text(separator='\n')
                
                # Clean whitespace
                lines = (line.strip() for line in text.splitlines())
                chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                text = '\n'.join(chunk for chunk in chunks if chunk)
                
                return title_text + text
            except Exception as e:
                return f"[Error processing HTML content: {str(e)}]"
        
        # Process spine documents in order
        processed_items = set()
        for item_id in book.spine:
            item = book.get_item_with_id(item_id[0] if isinstance(item_id, tuple) else item_id)
            if item and item.get_content() and item.media_type == 'application/xhtml+xml':
                processed_items.add(item.id)
                chapter_text = chapter_to_text(item.get_content().decode('utf-8'))
                if chapter_text:
                    content.append(f"--- Document: {item.get_name()} ---")
                    content.append(chapter_text)
                    content.append("-" * 40)
        
        # Process any HTML items not in spine but might contain important content
        for item in book.items:
            if (item.id not in processed_items and 
                item.media_type == 'application/xhtml+xml' and 
                not item.get_name().startswith('nav')):  # Skip navigation files
                chapter_text = chapter_to_text(item.get_content().decode('utf-8'))
                if chapter_text:
                    content.append(f"--- Additional Document: {item.get_name()} ---")
                    content.append(chapter_text)
                    content.append("-" * 40)
        
        return "\n".join(content)
    except Exception as e:
        return f"Error reading EPUB file: {str(e)}"

def read_rtf_file(file_path: str) -> str:
    """Extract text content from RTF files."""
    if not has_rtf_support:
        return "RTF support not available. Install required package with: pip install striprtf"
    
    try:
        # Read RTF content
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            rtf_text = file.read()
        
        # Convert RTF to plain text
        plain_text = striprtf.rtf_to_text(rtf_text)
        
        # Format output
        content = []
        content.append(f"--- RTF Document: {os.path.basename(file_path)} ---")
        content.append(f"Size: {os.path.getsize(file_path) / 1024:.2f} KB")
        content.append("-" * 40)
        content.append(plain_text)
        
        return "\n".join(content)
    except Exception as e:
        return f"Error reading RTF file: {str(e)}"

def extract_file_content(file_path: str) -> Dict[str, Any]:
    """Extract content from a file based on its extension."""
    if not os.path.exists(file_path):
        return {"success": False, "error": f"File not found: {file_path}", "content": ""}
    
    file_extension = os.path.splitext(file_path)[1].lower()
    content = ""
    
    try:
        # Handle different file types
        if file_extension in ['.txt', '.md', '.json', '.html', '.xml', '.log', '.py', '.js', '.css', '.java', '.ini', '.conf', '.cfg']:
            content = read_text_file(file_path)
        elif file_extension == '.pdf':
            content = read_pdf_file(file_path)
        elif file_extension == '.docx':
            content = read_docx_file(file_path)
        elif file_extension == '.xlsx':
            content = read_xlsx_file(file_path)
        elif file_extension == '.pptx':
            content = read_pptx_file(file_path)
        elif file_extension == '.csv':
            content = read_csv_file(file_path)
        elif file_extension == '.epub':
            content = read_epub_file(file_path)
        elif file_extension == '.rtf':
            content = read_rtf_file(file_path)
        else:
            # Try to read as text file first, then fall back to binary warning
            try:
                content = read_text_file(file_path)
            except Exception:
                return {
                    "success": False, 
                    "error": f"Unsupported file type: {file_extension}", 
                    "content": f"This file type ({file_extension}) is not directly supported."
                }
        
        return {"success": True, "content": content, "file_path": file_path, "file_type": file_extension}
    
    except Exception as e:
        return {"success": False, "error": str(e), "content": ""}

def summarize_content(content: str, max_length: int = 500) -> str:
    """Create a brief summary of the content if it's too long."""
    if len(content) <= max_length:
        return content
    
    # Simple summarization: first few characters and last few characters
    first_part = content[:max_length // 2]
    last_part = content[-(max_length // 2):]
    return f"{first_part}\n\n... [Content truncated, total length: {len(content)} characters] ...\n\n{last_part}"

def read_file(path: str, summarize: bool = False, max_summary_length: int = 500) -> Dict[str, Any]:
    """
    Reads and returns the contents of the file at 'path' with appropriate handling per file type.
    
    Args:
        path (str): Path to the file to read
        summarize (bool): Whether to summarize very large content
        max_summary_length (int): Maximum length for summary if summarizing
        
    Returns:
        Dict with keys:
        - success (bool): Whether the read was successful
        - content (str): The file content, possibly summarized
        - file_path (str): Original file path
        - file_type (str): File extension
        - error (str, optional): Error message if success is False
    """
    result = extract_file_content(path)
    
    # Summarize very large content if requested
    if summarize and result["success"] and len(result["content"]) > max_summary_length:
        result["content"] = summarize_content(result["content"], max_summary_length)
        result["summarized"] = True
    
    return result

def print_output(result: Dict[str, Any], output_format: str = "text") -> None:
    """Print the output in the specified format."""
    if output_format == "json":
        print(json.dumps(result, indent=2))
    else:
        if result["success"]:
            print(f"File: {result['file_path']} ({result['file_type']})")
            print("-" * 40)
            print(result["content"])
        else:
            print(f"Error: {result['error']}")

def save_to_file(content: str, original_path: str) -> str:
    """Save content to a text file and return the path."""
    # Create output filename based on original
    base_name = os.path.splitext(os.path.basename(original_path))[0]
    output_path = f"{base_name}_extracted.txt"
    
    # Save content to file
    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)
        return output_path
    except Exception as e:
        print(f"Error saving to file: {str(e)}")
        return None

def interactive_mode():
    """Run the file reader in interactive mode."""
    print("=" * 60)
    print("       Advanced File Content Extractor")
    print("=" * 60)
    print("Supported file types:")
    print("  • Text files (.txt, .md, .log, etc.)")
    print("  • PDF documents (.pdf)")
    print("  • Microsoft Word documents (.docx)")
    print("  • Microsoft Excel spreadsheets (.xlsx)")
    print("  • Microsoft PowerPoint presentations (.pptx)")
    print("  • CSV files (.csv)")
    print("  • EPUB e-books (.epub)")
    print("  • RTF documents (.rtf)")
    
    # Show which enhanced modules are available
    print("\nEnhanced features available:")
    if has_pymupdf:
        print("  ✓ Advanced PDF processing (PyMuPDF)")
    else:
        print("  ✗ Advanced PDF processing (install pymupdf for better results)")
    
    if has_epub_support:
        print("  ✓ EPUB e-book support")
    else:
        print("  ✗ EPUB support (install ebooklib and beautifulsoup4)")
    
    if has_rtf_support:
        print("  ✓ RTF document support")
    else:
        print("  ✗ RTF support (install striprtf)")
    
    if has_pil:
        print("  ✓ Enhanced image analysis")
    else:
        print("  ✗ Enhanced image analysis (install Pillow)")
    
    print("=" * 60)
    
    # Ask for file path
    file_path = input("Enter the path to the file you want to read: ").strip()
    
    # Check if file exists
    if not os.path.exists(file_path):
        print(f"Error: File '{file_path}' does not exist!")
        return
    
    # Options
    print("\nOutput options:")
    print("1. Display in terminal")
    print("2. Save to text file")
    print("3. Both display and save")
    
    while True:
        try:
            output_choice = int(input("Enter your choice (1-3): "))
            if 1 <= output_choice <= 3:
                break
            else:
                print("Please enter a number between 1 and 3.")
        except ValueError:
            print("Please enter a valid number.")
    
    # Summarize option for large files
    summarize = input("\nSummarize large content? (y/n): ").lower().startswith('y')
    max_length = 1000
    if summarize:
        try:
            max_length = int(input("Maximum length for summary (default: 1000): ") or "1000")
        except ValueError:
            print("Using default length of 1000 characters.")
    
    # Process file
    print(f"\nProcessing {file_path}...")
    start_time = time.time()
    result = read_file(file_path, summarize, max_length)
    processing_time = time.time() - start_time
    
    # Output
    if not result["success"]:
        print(f"\nError: {result['error']}")
        return
    
    file_size = os.path.getsize(file_path) / 1024  # KB
    print(f"\nFile processed successfully in {processing_time:.2f} seconds.")
    print(f"File size: {file_size:.2f} KB")
    print(f"File type: {result['file_type']}")
    content_length = len(result["content"])
    print(f"Extracted content length: {content_length} characters")
    
    # Display content
    if output_choice in [1, 3]:
        print("\n" + "=" * 60)
        print(f"CONTENT OF {os.path.basename(file_path)}:")
        print("=" * 60)
        print(result["content"])
        print("=" * 60)
    
    # Save to file
    if output_choice in [2, 3]:
        output_path = save_to_file(result["content"], file_path)
        if output_path:
            print(f"\nContent saved to: {output_path}")
    
    print("\nDone!")

if __name__ == '__main__':
    # Check if arguments were provided
    if len(sys.argv) > 1:
        # Command line mode
        import argparse
        parser = argparse.ArgumentParser(description='Read and extract content from various file types.')
        parser.add_argument('file_path', help='Path to the file to read')
        parser.add_argument('--format', choices=['text', 'json'], default='text', help='Output format')
        parser.add_argument('--summarize', action='store_true', help='Summarize large content')
        parser.add_argument('--max-length', type=int, default=1000, help='Maximum length for summary')
        parser.add_argument('--save', action='store_true', help='Save output to text file')
        
        args = parser.parse_args()
        
        # Process the file
        result = read_file(args.file_path, args.summarize, args.max_length)
        print_output(result, args.format)
        
        # Save to file if requested
        if args.save and result["success"]:
            output_path = save_to_file(result["content"], args.file_path)
            if output_path:
                print(f"Content saved to: {output_path}")
    else:
        # Interactive mode
        interactive_mode()