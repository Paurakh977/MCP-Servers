"""Office document reader module for Word, Excel and PowerPoint files."""
import os
import json
from datetime import datetime
from typing import Dict, Any, List, Optional, Union

# For image detection and properties
try:
    from PIL import Image
    has_pil = True
except ImportError:
    print("Pillow not installed. For better image analysis: pip install Pillow")
    has_pil = False

# For Microsoft Office documents
try:
    import docx  # for .docx files
    from docx.oxml.ns import qn  # for accessing embedded objects
    from docx.oxml import OxmlElement
except ImportError:
    print("python-docx not installed. To read Word files: pip install python-docx")

try:
    import openpyxl  # for .xlsx files
    from openpyxl.drawing.image import Image as XlsxImage
    from openpyxl.utils import get_column_letter
except ImportError:
    print("openpyxl not installed. To read Excel files: pip install openpyxl")

try:
    from pptx import Presentation  # for .pptx files
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("python-pptx not installed. To read PowerPoint files: pip install python-pptx")

# Function definitions will go here (added in separate edit)
def read_docx_file(file_path: str) -> str:
    """Extract text and identify objects from Microsoft Word (.docx) files in sequential order."""
    try:
        doc = docx.Document(file_path)
        full_text = []
        
        # Document properties first
        try:
            full_text.append("--- Document Properties ---")
            core_props = doc.core_properties
            if hasattr(core_props, 'title') and core_props.title:
                full_text.append(f"Title: {core_props.title}")
            if hasattr(core_props, 'author') and core_props.author:
                full_text.append(f"Author: {core_props.author}")
            if hasattr(core_props, 'created') and core_props.created:
                full_text.append(f"Created: {core_props.created}")
            if hasattr(core_props, 'modified') and core_props.modified:
                full_text.append(f"Modified: {core_props.modified}")
            if hasattr(core_props, 'comments') and core_props.comments:
                full_text.append(f"Comments: {core_props.comments}")
            if hasattr(core_props, 'category') and core_props.category:
                full_text.append(f"Category: {core_props.category}")
            if hasattr(core_props, 'subject') and core_props.subject:
                full_text.append(f"Subject: {core_props.subject}")
            if hasattr(core_props, 'keywords') and core_props.keywords:
                full_text.append(f"Keywords: {core_props.keywords}")
            full_text.append("-" * 40)
        except:
            pass  # Ignore if properties can't be accessed
        
        # Document statistics
        full_text.append("--- Document Statistics ---")
        full_text.append(f"Paragraphs: {len(doc.paragraphs)}")
        full_text.append(f"Sections: {len(doc.sections)}")
        full_text.append(f"Tables: {len(doc.tables)}")
        
        # Detect images and other embedded objects
        try:
            # Scan for embedded objects by looking at the underlying XML
            image_count = 0
            chart_count = 0
            shape_count = 0
            drawing_count = 0
            
            for rel in doc.part.rels.values():
                if rel.reltype == 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image':
                    image_count += 1
                elif rel.reltype == 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart':
                    chart_count += 1
                elif rel.reltype == 'http://schemas.microsoft.com/office/2007/relationships/shape':
                    shape_count += 1
                elif 'drawing' in rel.reltype:
                    drawing_count += 1
            
            if image_count > 0:
                full_text.append(f"Images: {image_count}")
                
                # Get more detailed image information if possible
                try:
                    image_details = []
                    for rel_id, rel in doc.part.rels.items():
                        if rel.reltype == 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image':
                            try:
                                # Try to get image dimensions
                                image_part = rel.target_part
                                if has_pil and image_part and hasattr(image_part, 'blob'):
                                    from io import BytesIO
                                    img = Image.open(BytesIO(image_part.blob))
                                    image_details.append(f"  Image {rel_id}: {img.format} {img.width}x{img.height}")
                                else:
                                    image_details.append(f"  Image {rel_id}")
                            except:
                                image_details.append(f"  Image {rel_id}")
                    
                    if image_details:
                        full_text.append("Image details:")
                        full_text.extend(image_details)
                except:
                    pass  # Skip if detailed image info can't be obtained
            
            if chart_count > 0:
                full_text.append(f"Charts: {chart_count}")
            if shape_count > 0:
                full_text.append(f"Shapes: {shape_count}")
            if drawing_count > 0:
                full_text.append(f"Drawings: {drawing_count}")
            
            full_text.append("-" * 40)
        except:
            full_text.append("-" * 40)  # Add separator even if object detection fails
        
        # Main content extraction - preserving sequence
        full_text.append("--- Document Content (in sequential order) ---")
        
        # We need to track all content elements (paragraphs, tables, images) and their positions
        content_elements = []
        
        # Track paragraphs and their sequence
        for i, para in enumerate(doc.paragraphs):
            # Skip empty paragraphs
            if not para.text.strip():
                continue
                
            # Check for headings
            heading_level = 0
            if para.style and para.style.name.startswith('Heading'):
                try:
                    heading_level = int(para.style.name.replace('Heading', ''))
                except ValueError:
                    heading_level = 0
            
            # Check for images or other objects within this paragraph
            has_objects = False
            has_drawings = False
            drawing_texts = []
            
            for run in para.runs:
                # Check for embedded images
                if run.element.findall('.//'+qn('w:drawing')) or run.element.findall('.//'+qn('w:pict')):
                    has_objects = True
                    
                    # Try to extract text from drawing elements (text boxes, etc.)
                    for drawing in run.element.findall('.//'+qn('w:drawing')):
                        has_drawings = True
                        # Try to extract any text from the drawing
                        try:
                            # Look for text elements in the drawing
                            texts = drawing.findall('.//' + qn('w:t'))
                            if texts:
                                drawing_text = ' '.join([t.text for t in texts if t.text])
                                if drawing_text.strip():
                                    drawing_texts.append(drawing_text.strip())
                        except:
                            pass
            
            # Create the paragraph text with appropriate heading level
            if heading_level > 0:
                para_text = f"{'#' * heading_level} {para.text}"
            else:
                para_text = para.text
            
            # Add information about embedded objects/drawings
            if has_objects:
                if drawing_texts:
                    para_text += f" [Contains drawing/image with text: {'; '.join(drawing_texts)}]"
                else:
                    para_text += " [Contains embedded object(s)]"
            
            # Store the paragraph with its document position
            content_elements.append({
                'type': 'paragraph',
                'position': i,
                'content': para_text,
                'has_objects': has_objects,
                'has_drawings': has_drawings,
                'drawing_texts': drawing_texts
            })
        
        # Track tables and their sequence
        for i, table in enumerate(doc.tables):
            # Find the paragraph immediately before the table to determine position
            table_position = -1
            
            # This is an approximation as the exact position might be harder to determine
            # We'll use XML to try to find where this table appears
            try:
                # Get all paragraphs and tables in order from the body
                body = doc._body._body
                elements = body.findall('.//')
                
                # Walk through elements to find this table's position
                current_position = 0
                for elem in elements:
                    if elem.tag.endswith('tbl'):  # Table element
                        if current_position == i:  # If this is our target table
                            table_position = current_position
                            break
                        current_position += 1
                    elif elem.tag.endswith('p'):  # Paragraph element
                        current_position += 1
            except:
                # If we can't determine position, just use the table index
                table_position = i * 1000  # Arbitrary large value to put tables after paragraphs
            
            # Format table data
            table_data = []
            col_widths = []
            
            # Process the table - similar to previous implementation
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    # Combine all text in the cell, handling paragraphs
                    cell_text = '\n'.join([p.text for p in cell.paragraphs if p.text.strip()])
                    
                    # Check for embedded objects in the cell
                    has_objects = False
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            if run.element.findall('.//'+qn('w:drawing')) or run.element.findall('.//'+qn('w:pict')):
                                has_objects = True
                                break
                    
                    if has_objects:
                        cell_text += " [Contains embedded object(s)]"
                        
                    # Clean up the cell text
                    cell_text = cell_text.strip()
                    row_data.append(cell_text)
                
                # Update column widths
                for i, cell_text in enumerate(row_data):
                    # For multiline content, get the max line width
                    lines = cell_text.split('\n')
                    max_line_width = max([len(line) for line in lines]) if lines else 0
                    
                    while len(col_widths) <= i:
                        col_widths.append(0)
                    col_widths[i] = max(col_widths[i], max_line_width)
                
                table_data.append(row_data)
            
            # Store the formatted table data
            content_elements.append({
                'type': 'table',
                'position': table_position,
                'content': table_data,
                'col_widths': col_widths
            })
        
        # Sort content elements by their position to maintain document flow
        content_elements.sort(key=lambda x: x['position'])
        
        # Now process each element in order
        for element in content_elements:
            if element['type'] == 'paragraph':
                full_text.append(element['content'])
            elif element['type'] == 'table':
                table_data = element['content']
                col_widths = element['col_widths']
                
                full_text.append("\n--- Table ---")
                
                # Format and output the table
                if table_data:
                    # Header row
                    header_row = table_data[0]
                    header_line = []
                    for i, cell in enumerate(header_row):
                        width = min(col_widths[i], 30)  # Limit width to reasonable size
                        header_line.append(cell.ljust(width))
                    full_text.append("| " + " | ".join(header_line) + " |")
                    
                    # Separator row
                    separator = []
                    for i in range(len(header_row)):
                        width = min(col_widths[i], 30)
                        separator.append("-" * width)
                    full_text.append("| " + " | ".join(separator) + " |")
                    
                    # Data rows
                    for row_idx, row_data in enumerate(table_data):
                        if row_idx == 0:  # Skip header row as we already displayed it
                            continue
                        
                        # Handle multiline cells
                        row_lines = []
                        max_lines = 1
                        
                        for i, cell in enumerate(row_data):
                            if not cell:
                                row_lines.append([""])
                                continue
                                
                            lines = cell.split('\n')
                            row_lines.append(lines)
                            max_lines = max(max_lines, len(lines))
                        
                        # Now create and output each line of the row
                        for line_idx in range(max_lines):
                            line_parts = []
                            for col_idx, cell_lines in enumerate(row_lines):
                                if line_idx < len(cell_lines):
                                    content = cell_lines[line_idx]
                                else:
                                    content = ""
                                    
                                width = min(col_widths[col_idx], 30)
                                line_parts.append(content.ljust(width))
                            
                            full_text.append("| " + " | ".join(line_parts) + " |")
                
                full_text.append("--- End Table ---\n")
        
        # Check for image-only or drawing-only objects not attached to paragraphs
        # This requires special handling at the XML level
        try:
            body = doc._body._body
            
            # Function to extract text from a drawing object
            def extract_drawing_text(drawing_elem):
                drawing_texts = []
                try:
                    texts = drawing_elem.findall('.//' + qn('w:t'))
                    if texts:
                        drawing_text = ' '.join([t.text for t in texts if t.text])
                        if drawing_text.strip():
                            drawing_texts.append(drawing_text.strip())
                except:
                    pass
                return drawing_texts
            
            # Find standalone drawing objects
            standalone_drawings = []
            for elem in body.findall('.//'+qn('w:drawing')):
                # Check if this drawing is directly in the body, not inside a paragraph we already processed
                parent = elem.getparent()
                if parent is not None and parent.tag.endswith('r'):  # Inside a run
                    run_parent = parent.getparent()
                    if run_parent is not None and run_parent.tag.endswith('p'):  # Inside a paragraph
                        # This might be a standalone drawing, check if we've already captured it
                        # This logic is approximate - ideally we'd match exact elements
                        drawing_texts = extract_drawing_text(elem)
                        if drawing_texts:
                            # Check if this text is in any of our existing paragraphs with drawings
                            found = False
                            for e in content_elements:
                                if e['type'] == 'paragraph' and e['has_drawings']:
                                    for dt in e['drawing_texts']:
                                        if dt in drawing_texts:
                                            found = True
                                            break
                            
                            if not found:
                                standalone_drawings.append({
                                    'type': 'drawing',
                                    'texts': drawing_texts
                                })
            
            # Add any standalone drawings we found
            if standalone_drawings:
                full_text.append("\n--- Standalone Drawings/Text Boxes ---")
                for i, drawing in enumerate(standalone_drawings):
                    full_text.append(f"Drawing/Text Box {i+1}:")
                    for text in drawing['texts']:
                        full_text.append(f"  {text}")
                full_text.append("--- End Drawings ---\n")
                
        except:
            pass  # Skip if we can't process standalone drawings
        
        return "\n".join(full_text)
    except Exception as e:
        return f"Error reading DOCX file: {str(e)}"

def read_xlsx_file(file_path: str, sheet_name: str = None, cell_range: str = None) -> str:
    """Extract data from Microsoft Excel (.xlsx) files with optimized output.
    
    Args:
        file_path: Path to the Excel file
        sheet_name: Optional specific sheet to read. If None, returns file info only
        cell_range: Optional cell range to read (e.g. 'A1:D10'). If None, reads all cells
        
    Returns:
        JSON string containing file info and optionally sheet data
    """
    try:
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        
        # Basic file info
        file_info = {
            "file": os.path.basename(file_path),
            "sheets": []
        }
        
        # Get info for all sheets
        for ws in workbook.worksheets:
            # Get sheet dimensions
            min_col, min_row, max_col, max_row = openpyxl.utils.range_boundaries(ws.calculate_dimension())
            
            # Get column headers (first row)
            columns = []
            column_refs = []
            for col in range(min_col, max_col + 1):
                cell = ws.cell(min_row, col)
                header = cell.value if cell.value is not None else f"Column {get_column_letter(col)}"
                columns.append(header)
                column_refs.append(get_column_letter(col))
            
            sheet_info = {
                "name": ws.title,
                "dimensions": ws.calculate_dimension(),
                "row_count": ws.max_row,
                "column_count": max_col - min_col + 1,
                "columns": columns,
                "column_refs": column_refs
            }
            file_info["sheets"].append(sheet_info)
        
        # If no specific sheet requested, return file info only
        if not sheet_name:
            return json.dumps(file_info)
        
        # Validate requested sheet exists
        if sheet_name not in workbook.sheetnames:
            return json.dumps({
                "error": f"Sheet '{sheet_name}' not found in workbook"
            })
        
        # Get the requested sheet
        sheet = workbook[sheet_name]
        
        # Parse cell range if provided
        if cell_range:
            try:
                min_col, min_row, max_col, max_row = openpyxl.utils.range_boundaries(cell_range)
            except ValueError:
                return json.dumps({
                    "error": f"Invalid cell range format: {cell_range}"
                })
        else:
            # Use full sheet range
            min_col, min_row, max_col, max_row = openpyxl.utils.range_boundaries(sheet.calculate_dimension())
        
        # Get column headers (first row)
        columns = []
        for col in range(min_col, max_col + 1):
            cell = sheet.cell(min_row, col)
            columns.append(cell.value if cell.value is not None else "")
        
        # Collect non-empty rows
        records = []
        for row in range(min_row + 1, max_row + 1):
            values = []
            row_has_data = False
            
            for col in range(min_col, max_col + 1):
                cell = sheet.cell(row, col)
                value = cell.value
                
                # Convert datetime objects to ISO format
                if isinstance(value, datetime.datetime):
                    value = value.isoformat()
                
                values.append(value)
                if value is not None and value != "":
                    row_has_data = True
            
            if row_has_data:
                records.append({
                    "row": row,
                    "values": values
                })
        
        # Count charts and images
        chart_count = len([drawing for drawing in sheet._charts])
        image_count = len([drawing for drawing in sheet._images])
        
        # Prepare sheet data
        sheet_data = {
            "sheet_name": sheet_name,
            "dimensions": cell_range or sheet.calculate_dimension(),
            "non_empty_cells": len(records),
            "charts": chart_count,
            "images": image_count,
            "columns": columns,
            "records": records
        }
        
        # Add sheet data to response
        file_info["sheet"] = [sheet_data]
        
        return json.dumps(file_info)
        
    except Exception as e:
        return json.dumps({
            "error": f"Failed to read Excel file: {str(e)}"
        })

def read_pptx_file(file_path: str) -> str:
    """Extract text and identify objects from Microsoft PowerPoint (.pptx) files in sequential order."""
    try:
        prs = Presentation(file_path)
        full_text = []
        
        # Document properties
        full_text.append("--- Presentation Properties ---")
        full_text.append(f"Filename: {os.path.basename(file_path)}")
        full_text.append(f"Number of slides: {len(prs.slides)}")
        
        # Try to get core properties
        try:
            core_props = prs.core_properties
            if hasattr(core_props, 'title') and core_props.title:
                full_text.append(f"Title: {core_props.title}")
            if hasattr(core_props, 'author') and core_props.author:
                full_text.append(f"Author: {core_props.author}")
            if hasattr(core_props, 'subject') and core_props.subject:
                full_text.append(f"Subject: {core_props.subject}")
            if hasattr(core_props, 'keywords') and core_props.keywords:
                full_text.append(f"Keywords: {core_props.keywords}")
            if hasattr(core_props, 'created') and core_props.created:
                full_text.append(f"Created: {core_props.created}")
            if hasattr(core_props, 'modified') and core_props.modified:
                full_text.append(f"Modified: {core_props.modified}")
        except:
            pass  # Skip if properties can't be accessed
        
        full_text.append("-" * 40)
        
        # Slide master and layouts info
        try:
            slide_masters = prs.slide_masters
            full_text.append(f"Slide masters: {len(slide_masters)}")
            
            total_layouts = 0
            for master in slide_masters:
                total_layouts += len(master.slide_layouts)
            
            full_text.append(f"Slide layouts: {total_layouts}")
            full_text.append("-" * 40)
        except:
            full_text.append("-" * 40)
        
        # Process each slide
        for i, slide in enumerate(prs.slides):
            slide_content = [f"--- Slide {i+1} ---"]
            
            # Extract slide layout type if available
            try:
                layout_name = slide.slide_layout.name
                slide_content.append(f"Layout: {layout_name}")
            except:
                pass
            
            # Get slide dimensions
            try:
                if hasattr(prs, 'slide_width') and hasattr(prs, 'slide_height'):
                    slide_width = prs.slide_width / 914400  # Convert EMU to inches
                    slide_height = prs.slide_height / 914400
                    slide_content.append(f"Dimensions: {slide_width:.2f}\" × {slide_height:.2f}\"")
            except:
                pass
                
            # Extract text from slide title separately
            title_text = ""
            if slide.shapes.title:
                title_text = slide.shapes.title.text
                slide_content.append(f"Title: {title_text}")
            
            # Count shape types for summary
            shape_counts = {
                'text_boxes': 0,
                'pictures': 0,
                'charts': 0,
                'tables': 0,
                'diagrams': 0,
                'videos': 0,
                'other_shapes': 0
            }
            
            # Collect all shapes with their position information for sequential ordering
            slide_elements = []
            
            # Process all shapes to extract content and position
            for shape_idx, shape in enumerate(slide.shapes):
                # Skip the title as we've already handled it
                if shape == slide.shapes.title:
                    continue
                
                # Get shape position
                try:
                    top = shape.top / 914400  # EMU to inches
                    left = shape.left / 914400
                except:
                    # If we can't get position, use index as proxy for position
                    top = shape_idx * 10
                    left = 0
                
                # Check type and extract content
                if shape.has_text_frame:
                    if shape.text.strip():
                        shape_counts['text_boxes'] += 1
                        
                        # Store text with paragraph formatting preserved
                        shape_text = []
                        
                        # Preserve text formatting - check paragraphs and runs
                        for para in shape.text_frame.paragraphs:
                            para_text = []
                            for run in para.runs:
                                text = run.text.strip()
                                if text:
                                    # Check for formatting
                                    if hasattr(run, 'font') and hasattr(run.font, 'bold') and run.font.bold:
                                        text = f"**{text}**"
                                    if hasattr(run, 'font') and hasattr(run.font, 'italic') and run.font.italic:
                                        text = f"*{text}*"
                                    if hasattr(run, 'font') and hasattr(run.font, 'underline') and run.font.underline:
                                        text = f"_{text}_"
                                    
                                    para_text.append(text)
                            
                            if para_text:
                                shape_text.append(" ".join(para_text))
                        
                        # Add the text content with formatting preserved
                        text_content = "\n".join(shape_text)
                        
                        slide_elements.append({
                            'type': 'text',
                            'top': top,
                            'left': left,
                            'content': text_content
                        })
                
                elif shape.has_table:
                    shape_counts['tables'] += 1
                    
                    # Process table data with column widths
                    table_data = []
                    col_widths = []
                    
                    # Extract table content
                    for row in shape.table.rows:
                        row_data = []
                        for cell in row.cells:
                            cell_text = cell.text.strip().replace('\n', ' ')
                            row_data.append(cell_text)
                        
                        # Update column widths
                        for i, text in enumerate(row_data):
                            while len(col_widths) <= i:
                                col_widths.append(0)
                            col_widths[i] = max(col_widths[i], len(text))
                        
                        table_data.append(row_data)
                    
                    slide_elements.append({
                        'type': 'table',
                        'top': top,
                        'left': left,
                        'content': table_data,
                        'col_widths': col_widths
                    })
                
                elif shape.has_chart:
                    shape_counts['charts'] += 1
                    
                    # Get chart type and data if possible
                    try:
                        chart_type = str(shape.chart.chart_type).split('.')[-1]
                        chart_desc = f"Chart: {chart_type}"
                    except:
                        chart_desc = "Chart"
                    
                    slide_elements.append({
                        'type': 'chart',
                        'top': top,
                        'left': left,
                        'content': chart_desc
                    })
                
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    shape_counts['pictures'] += 1
                    
                    # Get picture details if possible
                    try:
                        if hasattr(shape, 'width') and hasattr(shape, 'height'):
                            pic_width = shape.width / 914400  # EMU to inches
                            pic_height = shape.height / 914400
                            pic_desc = f"Picture: {pic_width:.1f}\" × {pic_height:.1f}\""
                        else:
                            pic_desc = "Picture"
                    except:
                        pic_desc = "Picture"
                    
                    slide_elements.append({
                        'type': 'picture',
                        'top': top,
                        'left': left,
                        'content': pic_desc
                    })
                
                elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                    shape_counts['videos'] += 1
                    
                    slide_elements.append({
                        'type': 'media',
                        'top': top,
                        'left': left,
                        'content': "Video/Media"
                    })
                
                elif shape.shape_type == MSO_SHAPE_TYPE.DIAGRAM:
                    shape_counts['diagrams'] += 1
                    
                    slide_elements.append({
                        'type': 'diagram',
                        'top': top,
                        'left': left,
                        'content': "Diagram/SmartArt"
                    })
                
                else:
                    shape_counts['other_shapes'] += 1
                    
                    # Try to get text from drawing objects or other shapes
                    shape_text = ""
                    try:
                        if hasattr(shape, 'text'):
                            shape_text = shape.text
                    except:
                        pass
                    
                    slide_elements.append({
                        'type': 'shape',
                        'top': top,
                        'left': left,
                        'content': f"Shape{': ' + shape_text if shape_text else ''}"
                    })
            
            # Add shape summary
            shape_summary = []
            for shape_type, count in shape_counts.items():
                if count > 0:
                    shape_summary.append(f"{count} {shape_type.replace('_', ' ')}")
            
            if shape_summary:
                slide_content.append("Objects: " + ", ".join(shape_summary))
            
            # Sort elements by position (top to bottom, then left to right)
            slide_elements.sort(key=lambda x: (x['top'], x['left']))
            
            # Now process all elements in sequential order
            slide_content.append("\n--- Slide Content (in sequential order) ---")
            
            for element in slide_elements:
                element_type = element['type']
                
                if element_type == 'text':
                    slide_content.append(f"Text Box:")
                    # Add indentation to preserve the text box structure
                    indented_text = "\n".join(f"  {line}" for line in element['content'].split('\n'))
                    slide_content.append(indented_text)
                
                elif element_type == 'picture':
                    slide_content.append(f"[{element['content']}]")
                
                elif element_type == 'chart':
                    slide_content.append(f"[{element['content']}]")
                
                elif element_type == 'diagram':
                    slide_content.append(f"[{element['content']}]")
                
                elif element_type == 'media':
                    slide_content.append(f"[{element['content']}]")
                
                elif element_type == 'shape':
                    slide_content.append(f"[{element['content']}]")
                
                elif element_type == 'table':
                    table_data = element['content']
                    col_widths = element['col_widths']
                    
                    slide_content.append("--- Table ---")
                    
                    # Format and output the table
                    if table_data:
                        # Header row
                        header_row = table_data[0]
                        header_parts = []
                        for i, text in enumerate(header_row):
                            width = min(col_widths[i], 30)  # Cap width at 30 chars
                            header_parts.append(text.ljust(width))
                        slide_content.append("| " + " | ".join(header_parts) + " |")
                        
                        # Separator
                        separator_parts = []
                        for i in range(len(header_row)):
                            width = min(col_widths[i], 30)
                            separator_parts.append("-" * width)
                        slide_content.append("| " + " | ".join(separator_parts) + " |")
                        
                        # Data rows
                        for row_idx, row_data in enumerate(table_data):
                            if row_idx == 0:  # Skip header as we already displayed it
                                continue
                                
                            row_parts = []
                            for i, text in enumerate(row_data):
                                width = min(col_widths[i], 30)
                                row_parts.append(text.ljust(width))
                            slide_content.append("| " + " | ".join(row_parts) + " |")
                    
                    slide_content.append("--- End Table ---")
            
            # Check if slide has minimal content
            if len(slide_content) <= 3:  # Only slide number, maybe layout and title
                slide_content.append("[This slide appears to have no text content or contains only non-text elements]")
            
            full_text.append("\n".join(slide_content))
        
        return "\n\n".join(full_text)
    except Exception as e:
        return f"Error reading PPTX file: {str(e)}" 